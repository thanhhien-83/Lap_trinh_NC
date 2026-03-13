# -*- coding: utf-8 -*-
"""
Script tạo file PowerPoint cho bài thuyết trình bảo vệ BTL
Đề tài: Xây dựng phần mềm quản lý sinh viên sử dụng cấu trúc mảng động và giao diện đồ họa
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============= COLOR PALETTE =============
PRIMARY = RGBColor(79, 70, 229)      # Indigo #4F46E5
PRIMARY_DARK = RGBColor(67, 56, 202) # #4338CA
SECONDARY = RGBColor(99, 102, 241)   # #6366F1
ACCENT = RGBColor(16, 185, 129)      # Green #10B981
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 30, 30)
DARK_TEXT = RGBColor(30, 41, 59)     # #1E293B
GRAY_TEXT = RGBColor(100, 116, 139)  # #64748B
LIGHT_BG = RGBColor(241, 245, 249)  # #F1F5F9
CARD_BG = RGBColor(248, 250, 252)   # #F8FAFC
ORANGE = RGBColor(249, 115, 22)     # #F97316
RED = RGBColor(239, 68, 68)         # #EF4444
YELLOW = RGBColor(234, 179, 8)      # #EAB308
BLUE = RGBColor(59, 130, 246)       # #3B82F6

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_gradient_bg(slide, color1=PRIMARY, color2=PRIMARY_DARK):
    """Add a solid background color to slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color1


def add_light_bg(slide):
    """Add light background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_decorative_bar(slide, top=0, height=Inches(0.08), color=PRIMARY):
    """Add a thin decorative bar at top"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bottom_bar(slide, text="", color=PRIMARY):
    """Add bottom bar with text"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = WHITE
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER


def add_slide_number(slide, num, total):
    """Add slide number"""
    txBox = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(1), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.color.rgb = WHITE
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.RIGHT


def add_section_header(slide, number, title, subtitle=""):
    """Add consistent section header to content slides"""
    add_light_bg(slide)
    add_decorative_bar(slide, top=Inches(0), height=Inches(0.06), color=PRIMARY)

    # Section number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.6), Inches(0.35), Inches(0.55), Inches(0.55)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.color.rgb = WHITE
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.35), Inches(0.3), Inches(10), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.color.rgb = DARK_TEXT
    p.font.size = Pt(28)
    p.font.bold = True

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.35), Inches(0.8), Inches(10), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.color.rgb = GRAY_TEXT
        p2.font.size = Pt(14)

    # Separator line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(226, 232, 240)
    line.line.fill.background()


def add_card(slide, left, top, width, height, title="", content="", icon="",
             title_color=DARK_TEXT, bg_color=CARD_BG, border_color=None, content_lines=None):
    """Add a card-style box"""
    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.color.rgb = RGBColor(226, 232, 240)
        card.line.width = Pt(1)

    y_offset = top + Inches(0.15)

    if icon:
        txIcon = slide.shapes.add_textbox(left + Inches(0.2), y_offset, Inches(0.4), Inches(0.35))
        tf = txIcon.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(18)

    if title:
        title_left = left + (Inches(0.6) if icon else Inches(0.2))
        txTitle = slide.shapes.add_textbox(title_left, y_offset, width - Inches(0.4), Inches(0.35))
        tf = txTitle.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.color.rgb = title_color
        p.font.size = Pt(14)
        p.font.bold = True
        y_offset += Inches(0.35)

    if content:
        txContent = slide.shapes.add_textbox(left + Inches(0.2), y_offset, width - Inches(0.4), height - Inches(0.5))
        tf = txContent.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.color.rgb = GRAY_TEXT
        p.font.size = Pt(12)
        p.space_before = Pt(2)

    if content_lines:
        txContent = slide.shapes.add_textbox(left + Inches(0.2), y_offset, width - Inches(0.4), height - Inches(0.5))
        tf = txContent.text_frame
        tf.word_wrap = True
        for i, line in enumerate(content_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.color.rgb = DARK_TEXT
            p.font.size = Pt(12)
            p.space_before = Pt(3)
            p.space_after = Pt(3)

    return card


def add_highlight_box(slide, left, top, width, height, text, color=PRIMARY, font_size=Pt(32)):
    """Add a highlight number box"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = WHITE
    p.font.size = font_size
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)
    return box


TOTAL_SLIDES = 14  # We'll track this

# ================================================================
# SLIDE 1: Title Slide
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_gradient_bg(slide, PRIMARY, PRIMARY_DARK)

# Decorative top line
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2), Inches(0.8), Inches(9.333), Inches(0.04)
)
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(255, 255, 255)
line.line.fill.background()
line.fill.fore_color.rgb = SECONDARY

# University / Course info
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.1), Inches(11.333), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "BÁO CÁO BÀI TẬP LỚN"
p.font.color.rgb = RGBColor(199, 210, 254)
p.font.size = Pt(18)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Môn học: LẬP TRÌNH NÂNG CAO"
p2.font.color.rgb = RGBColor(199, 210, 254)
p2.font.size = Pt(14)
p2.alignment = PP_ALIGN.CENTER

# Main title
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "XÂY DỰNG PHẦN MỀM QUẢN LÝ SINH VIÊN"
p.font.color.rgb = WHITE
p.font.size = Pt(36)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "SỬ DỤNG CẤU TRÚC MẢNG ĐỘNG VÀ GIAO DIỆN ĐỒ HỌA"
p2.font.color.rgb = RGBColor(199, 210, 254)
p2.font.size = Pt(24)
p2.font.bold = True
p2.alignment = PP_ALIGN.CENTER

# Decorative accent shape
accent = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.2), Inches(2.333), Inches(0.06)
)
accent.fill.solid()
accent.fill.fore_color.rgb = ACCENT
accent.line.fill.background()

# Tech badges
badges = ["HTML5", "CSS3", "JavaScript", "Chart.js", "LocalStorage"]
badge_start_x = Inches(3.0)
for i, badge_text in enumerate(badges):
    bx = badge_start_x + Inches(i * 1.6)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(4.6), Inches(1.4), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(99, 102, 241)
    badge.line.color.rgb = RGBColor(129, 140, 248)
    badge.line.width = Pt(1)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = badge_text
    p.font.color.rgb = WHITE
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

# Student info
txBox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Sinh viên thực hiện: [Điền tên sinh viên]"
p.font.color.rgb = RGBColor(199, 210, 254)
p.font.size = Pt(16)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "GVHD: [Điền tên giảng viên]  •  Năm học: 2025-2026"
p2.font.color.rgb = RGBColor(165, 180, 252)
p2.font.size = Pt(13)
p2.alignment = PP_ALIGN.CENTER

# Bottom line
line2 = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2), Inches(6.9), Inches(9.333), Inches(0.04)
)
line2.fill.solid()
line2.fill.fore_color.rgb = SECONDARY
line2.line.fill.background()


# ================================================================
# SLIDE 2: Nội dung trình bày (Agenda)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "📋", "NỘI DUNG TRÌNH BÀY", "Tổng quan các phần trong bài bảo vệ")
add_bottom_bar(slide, "Bài tập lớn - Lập trình nâng cao", PRIMARY)
add_slide_number(slide, 2, TOTAL_SLIDES)

agenda_items = [
    ("01", "Đặt vấn đề & Mục tiêu", "Giới thiệu bài toán và mục tiêu đề tài", PRIMARY),
    ("02", "Cơ sở lý thuyết", "Mảng động, QuickSort, Binary Search", SECONDARY),
    ("03", "Phân tích & Thiết kế", "Mô hình dữ liệu, kiến trúc hệ thống", BLUE),
    ("04", "Thuật toán & Độ phức tạp", "Chi tiết cài đặt và phân tích", ACCENT),
    ("05", "Demo ứng dụng", "Giao diện và các chức năng chính", ORANGE),
    ("06", "Kết luận & Hướng phát triển", "Tổng kết và đề xuất", RED),
]

for i, (num, title, desc, color) in enumerate(agenda_items):
    row = i // 3
    col = i % 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.5 + row * 2.6)

    # Number badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.55), Inches(0.55)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.color.rgb = WHITE
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Title
    txBox = slide.shapes.add_textbox(left + Inches(0.7), top, Inches(3.2), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.color.rgb = DARK_TEXT
    p.font.size = Pt(16)
    p.font.bold = True

    # Description
    txBox2 = slide.shapes.add_textbox(left + Inches(0.7), top + Inches(0.4), Inches(3.2), Inches(0.4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.color.rgb = GRAY_TEXT
    p2.font.size = Pt(12)

    # Underline accent
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + Inches(0.7), top + Inches(0.85), Inches(2.5), Inches(0.03)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = RGBColor(226, 232, 240)
    underline.line.fill.background()


# ================================================================
# SLIDE 3: Đặt vấn đề
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "1", "ĐẶT VẤN ĐỀ & MỤC TIÊU", "Bối cảnh và mục tiêu của đề tài")
add_bottom_bar(slide, "Phần 1: Đặt vấn đề & Mục tiêu", PRIMARY)
add_slide_number(slide, 3, TOTAL_SLIDES)

# Problem cards
add_card(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.5),
         title="⚡ Vấn đề thực tế", 
         content_lines=[
             "• Quản lý thông tin sinh viên là nhu cầu thiết yếu",
             "  trong môi trường giáo dục hiện đại",
             "• Cần hệ thống CRUD nhanh chóng, hiệu quả",
             "• Cần chức năng tìm kiếm, sắp xếp linh hoạt",
             "• Thống kê kết quả học tập trực quan",
             "• Giao diện thân thiện, dễ sử dụng",
         ],
         border_color=PRIMARY)

add_card(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(2.5),
         title="🎯 Mục tiêu đề tài",
         content_lines=[
             "• Xây dựng cấu trúc mảng động từ đầu",
             "  (không dùng thư viện có sẵn)",
             "• Cài đặt QuickSort - O(n log n)",
             "• Cài đặt Binary Search - O(log n)",
             "• Phát triển giao diện đồ họa hiện đại",
             "• Đảm bảo hiệu năng & khả năng mở rộng",
         ],
         border_color=ACCENT)

# Technology stack section
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "🛠 Công nghệ sử dụng"
p.font.color.rgb = DARK_TEXT
p.font.size = Pt(16)
p.font.bold = True

techs = [
    ("HTML5", "Cấu trúc trang web", PRIMARY),
    ("CSS3", "Giao diện & Animation", SECONDARY),
    ("JavaScript ES6+", "Logic nghiệp vụ & OOP", ACCENT),
    ("Chart.js", "Biểu đồ thống kê", ORANGE),
    ("LocalStorage", "Lưu trữ dữ liệu", BLUE),
]
for i, (tech, desc, color) in enumerate(techs):
    left = Inches(0.6 + i * 2.5)
    add_card(slide, left, Inches(4.8), Inches(2.3), Inches(1.5),
             title=tech, content=desc, title_color=color)


# ================================================================
# SLIDE 4: Cấu trúc Mảng Động
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "2", "CẤU TRÚC MẢNG ĐỘNG (Dynamic Array)", "Cơ sở lý thuyết - Cấu trúc dữ liệu chính")
add_bottom_bar(slide, "Phần 2: Cơ sở lý thuyết", PRIMARY)
add_slide_number(slide, 4, TOTAL_SLIDES)

# Definition card
add_card(slide, Inches(0.6), Inches(1.4), Inches(6.0), Inches(2.2),
         title="📦 Định nghĩa & Đặc điểm",
         content_lines=[
             "Mảng động cho phép lưu trữ phần tử liên tiếp",
             "với khả năng tự động thay đổi kích thước.",
             "",
             "✓ Truy cập ngẫu nhiên O(1)",
             "✓ Tự động mở rộng (capacity × 2)",
             "✓ Bộ nhớ liên tục → tối ưu cache",
         ],
         border_color=PRIMARY)

# Resize mechanism
add_card(slide, Inches(7.0), Inches(1.4), Inches(5.7), Inches(2.2),
         title="🔄 Cơ chế Resize",
         content_lines=[
             "Khi length >= capacity:",
             "  1. Tạo mảng mới: capacity × 2",
             "  2. Copy dữ liệu sang mảng mới",
             "  3. Giải phóng mảng cũ",
             "",
             "→ Amortized O(1) cho mỗi push()",
         ],
         border_color=ACCENT)

# Amortized Analysis
add_card(slide, Inches(0.6), Inches(3.9), Inches(5.5), Inches(2.8),
         title="📊 Amortized Analysis",
         content_lines=[
             "Thêm n phần tử vào mảng động:",
             "",
             "  Số lần resize: log₂(n)",
             "  Tổng số copy:  1+2+4+...+n = 2n−1",
             "  Chi phí TB:    (2n−1+n)/n ≈ 3 = O(1)",
             "",
             "→ Chi phí trung bình mỗi push() là O(1)!",
         ],
         border_color=BLUE)

# Complexity table
txBox = slide.shapes.add_textbox(Inches(6.5), Inches(3.9), Inches(6), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "📋 Bảng độ phức tạp DynamicArray"
p.font.color.rgb = DARK_TEXT
p.font.size = Pt(14)
p.font.bold = True

# Table
table_data = [
    ["Thao tác", "Thời gian", "Ghi chú"],
    ["push()", "O(1)*", "Amortized"],
    ["pop()", "O(1)", ""],
    ["get() / set()", "O(1)", "Random access"],
    ["insertAt()", "O(n)", "Dịch phần tử"],
    ["removeAt()", "O(n)", "Dịch phần tử"],
    ["find() / filter()", "O(n)", "Linear search"],
    ["sort()", "O(n log n)", "QuickSort"],
    ["binarySearch()", "O(log n)", "Đã sorted"],
]

table = slide.shapes.add_table(len(table_data), 3, Inches(6.5), Inches(4.35), Inches(6.1), Inches(2.3)).table
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(1.6)
table.columns[2].width = Inches(2.3)

for row_idx, row_data in enumerate(table_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        if row_idx == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY
        else:
            p.font.color.rgb = DARK_TEXT
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE


# ================================================================
# SLIDE 5: QuickSort Algorithm
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "2", "THUẬT TOÁN QUICKSORT", "Cơ sở lý thuyết - Thuật toán sắp xếp")
add_bottom_bar(slide, "Phần 2: Cơ sở lý thuyết", PRIMARY)
add_slide_number(slide, 5, TOTAL_SLIDES)

# Idea card
add_card(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.0),
         title="💡 Ý tưởng: Chia để trị (Divide & Conquer)",
         content_lines=[
             "1. Chọn phần tử chốt (pivot)",
             "2. Phân hoạch: chia mảng thành 2 phần",
             "   • Phần trái: các phần tử ≤ pivot",
             "   • Phần phải: các phần tử > pivot",
             "3. Đệ quy sắp xếp 2 phần con",
         ],
         border_color=PRIMARY)

# Pseudocode
add_card(slide, Inches(6.9), Inches(1.4), Inches(5.8), Inches(2.0),
         title="📝 Pseudocode",
         content_lines=[
             "QUICKSORT(arr, low, high):",
             "  if low < high:",
             "    pi = PARTITION(arr, low, high)",
             "    QUICKSORT(arr, low, pi - 1)",
             "    QUICKSORT(arr, pi + 1, high)",
         ],
         border_color=SECONDARY)

# Partition pseudocode
add_card(slide, Inches(0.6), Inches(3.7), Inches(5.8), Inches(2.7),
         title="📝 Hàm PARTITION",
         content_lines=[
             "PARTITION(arr, low, high):",
             "  pivot = arr[high]",
             "  i = low - 1",
             "  for j = low to high - 1:",
             "    if arr[j] <= pivot:",
             "      i = i + 1",
             "      swap(arr[i], arr[j])",
             "  swap(arr[i+1], arr[high])",
             "  return i + 1",
         ],
         border_color=BLUE)

# Complexity card
add_card(slide, Inches(6.9), Inches(3.7), Inches(5.8), Inches(2.7),
         title="⏱ Độ phức tạp",
         content_lines=[
             "",
             "  Trường hợp     | Thời gian    | Không gian",
             "  ─────────────────────────────────────",
             "  Tốt nhất         | O(n log n)   | O(log n)",
             "  Trung bình     | O(n log n)   | O(log n)",
             "  Xấu nhất        | O(n²)           | O(n)",
             "",
             "→ TB: O(n log n) - hiệu quả cho dữ liệu lớn",
         ],
         border_color=ACCENT)


# ================================================================
# SLIDE 6: Binary Search & Linear Search
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "2", "THUẬT TOÁN TÌM KIẾM", "Cơ sở lý thuyết - Tìm kiếm tuyến tính & nhị phân")
add_bottom_bar(slide, "Phần 2: Cơ sở lý thuyết", PRIMARY)
add_slide_number(slide, 6, TOTAL_SLIDES)

# Linear Search
add_card(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.5),
         title="🔍 Tìm kiếm tuyến tính (Linear Search)",
         content_lines=[
             "Ý tưởng: Duyệt từng phần tử từ đầu → cuối",
             "",
             "LINEAR_SEARCH(arr, target):",
             "  for i = 0 to length - 1:",
             "    if arr[i] == target: return i",
             "  return -1",
             "",
             "⏱ Độ phức tạp: O(n)",
         ],
         border_color=ORANGE)

# Binary Search
add_card(slide, Inches(6.9), Inches(1.4), Inches(5.8), Inches(2.5),
         title="⚡ Tìm kiếm nhị phân (Binary Search)",
         content_lines=[
             "Ý tưởng: Chia đôi mảng đã sắp xếp",
             "",
             "BINARY_SEARCH(arr, target):",
             "  left=0, right=length-1",
             "  while left <= right:",
             "    mid = (left+right)/2",
             "    if arr[mid] == target: return mid",
             "    else if arr[mid] < target: left=mid+1",
             "    else: right=mid-1",
         ],
         border_color=ACCENT)

# Comparison
add_card(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(2.3),
         title="📊 So sánh hiệu quả", border_color=PRIMARY)

comp_table_data = [
    ["Tiêu chí", "Linear Search", "Binary Search"],
    ["Độ phức tạp", "O(n)", "O(log n)"],
    ["Yêu cầu", "Không cần sắp xếp", "Mảng đã sắp xếp"],
    ["100 phần tử", "Tối đa 100 bước", "Tối đa 7 bước"],
    ["1000 phần tử", "Tối đa 1000 bước", "Tối đa 10 bước"],
    ["Ứng dụng", "Tìm theo nhiều tiêu chí", "Tìm nhanh theo mã SV"],
]

comp_table = slide.shapes.add_table(len(comp_table_data), 3, Inches(1.0), Inches(4.65), Inches(11.2), Inches(1.7)).table
comp_table.columns[0].width = Inches(2.8)
comp_table.columns[1].width = Inches(4.2)
comp_table.columns[2].width = Inches(4.2)

for row_idx, row_data in enumerate(comp_table_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = comp_table.cell(row_idx, col_idx)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        if row_idx == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY
        else:
            p.font.color.rgb = DARK_TEXT
            if col_idx == 2:
                p.font.color.rgb = ACCENT
                p.font.bold = True
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)


# ================================================================
# SLIDE 7: Phân tích bài toán
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "3", "PHÂN TÍCH BÀI TOÁN", "Yêu cầu chức năng và phi chức năng")
add_bottom_bar(slide, "Phần 3: Phân tích & Thiết kế", PRIMARY)
add_slide_number(slide, 7, TOTAL_SLIDES)

# Functional requirements
func_items = [
    ("CRUD", "Thêm, Xem, Sửa, Xóa\nsinh viên"),
    ("Tìm kiếm", "Tìm theo mã SV,\nhọ tên, lớp, xếp loại"),
    ("Sắp xếp", "Theo tên, mã SV,\nGPA, lớp (A-Z, Z-A)"),
    ("Thống kê", "Phân bố xếp loại,\ntheo lớp, biểu đồ"),
]

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(6), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "📋 Yêu cầu chức năng"
p.font.color.rgb = DARK_TEXT
p.font.size = Pt(16)
p.font.bold = True

colors = [PRIMARY, ACCENT, BLUE, ORANGE]
for i, (title, desc) in enumerate(func_items):
    left = Inches(0.6 + i * 3.1)
    add_card(slide, left, Inches(1.85), Inches(2.8), Inches(1.8),
             title=title, content=desc, title_color=colors[i],
             border_color=colors[i])

# Non-functional requirements
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.9), Inches(6), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "⚙️ Yêu cầu phi chức năng"
p.font.color.rgb = DARK_TEXT
p.font.size = Pt(16)
p.font.bold = True

nf_items = [
    ("Hiệu năng", "Các thao tác phải nhanh\nvới < 1000 sinh viên"),
    ("Khả dụng", "Giao diện trực quan,\ndễ sử dụng"),
    ("Lưu trữ", "Dữ liệu lưu vào\nLocalStorage trình duyệt"),
    ("Tương thích", "Hoạt động trên các\ntrình duyệt web hiện đại"),
]

for i, (title, desc) in enumerate(nf_items):
    left = Inches(0.6 + i * 3.1)
    add_card(slide, left, Inches(4.4), Inches(2.8), Inches(1.5),
             title=title, content=desc, title_color=GRAY_TEXT)

# Data model summary
add_card(slide, Inches(0.6), Inches(6.05), Inches(12), Inches(0.75),
         title="📊 Mô hình: Student { id, name, dob, gender, className, email, phone, address, scores(math, lit, eng), GPA = (Toán + Văn + Anh)/3 }",
         title_color=PRIMARY)


# ================================================================
# SLIDE 8: Kiến trúc hệ thống
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "3", "KIẾN TRÚC HỆ THỐNG", "Thiết kế 3 tầng và cấu trúc file")
add_bottom_bar(slide, "Phần 3: Phân tích & Thiết kế", PRIMARY)
add_slide_number(slide, 8, TOTAL_SLIDES)

# Layer 1: Presentation
layer1 = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.4), Inches(10.3), Inches(1.2)
)
layer1.fill.solid()
layer1.fill.fore_color.rgb = RGBColor(238, 242, 255)
layer1.line.color.rgb = PRIMARY
layer1.line.width = Pt(2)

txBox = slide.shapes.add_textbox(Inches(1.8), Inches(1.45), Inches(3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "PRESENTATION LAYER (Tầng giao diện)"
p.font.color.rgb = PRIMARY
p.font.size = Pt(13)
p.font.bold = True

pres_items = ["index.html", "styles.css", "Chart.js"]
for i, item in enumerate(pres_items):
    bx = Inches(2.0 + i * 3.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(1.9), Inches(2.5), Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(199, 210, 254)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

# Arrow 1
arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.4), Inches(2.65), Inches(0.5), Inches(0.4))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = GRAY_TEXT
arrow1.line.fill.background()

# Layer 2: Business Logic
layer2 = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.1), Inches(10.3), Inches(1.7)
)
layer2.fill.solid()
layer2.fill.fore_color.rgb = RGBColor(236, 253, 245)
layer2.line.color.rgb = ACCENT
layer2.line.width = Pt(2)

txBox = slide.shapes.add_textbox(Inches(1.8), Inches(3.15), Inches(4), Inches(0.35))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "BUSINESS LOGIC LAYER (Tầng xử lý nghiệp vụ)"
p.font.color.rgb = ACCENT
p.font.size = Pt(13)
p.font.bold = True

biz_items = [
    ("DieuKhienGiaoDien.js", "UI Controller"),
    ("QuanLySinhVien.js", "Student Manager"),
    ("ThuatToanSapXep.js", "QuickSort"),
    ("ThuatToanTimKiem.js", "Search"),
    ("ThuatToanThongKe.js", "Statistics"),
]
for i, (name, desc) in enumerate(biz_items):
    bx = Inches(1.8 + i * 2.0)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(3.55), Inches(1.8), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(167, 243, 208)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(9)
    p.font.color.rgb = ACCENT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY_TEXT
    p2.alignment = PP_ALIGN.CENTER

# Arrow 2
arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.4), Inches(4.85), Inches(0.5), Inches(0.4))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = GRAY_TEXT
arrow2.line.fill.background()

# Layer 3: Data
layer3 = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.3), Inches(10.3), Inches(1.3)
)
layer3.fill.solid()
layer3.fill.fore_color.rgb = RGBColor(255, 247, 237)
layer3.line.color.rgb = ORANGE
layer3.line.width = Pt(2)

txBox = slide.shapes.add_textbox(Inches(1.8), Inches(5.35), Inches(3), Inches(0.35))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "DATA LAYER (Tầng dữ liệu)"
p.font.color.rgb = ORANGE
p.font.size = Pt(13)
p.font.bold = True

data_items = [
    ("MangDong.js", "Dynamic Array"),
    ("SinhVien.js", "Student Model"),
    ("CauTrucMangDong.js", "Theory docs"),
    ("LocalStorage", "Persistence"),
]
for i, (name, desc) in enumerate(data_items):
    bx = Inches(2.0 + i * 2.4)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(5.75), Inches(2.0), Inches(0.65))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(253, 186, 116)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{name} ({desc})"
    p.font.size = Pt(10)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER


# ================================================================
# SLIDE 9: Thiết kế lớp chi tiết
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "3", "THIẾT KẾ CÁC LỚP (OOP)", "Class Diagram - Các lớp chính trong hệ thống")
add_bottom_bar(slide, "Phần 3: Phân tích & Thiết kế", PRIMARY)
add_slide_number(slide, 9, TOTAL_SLIDES)

# DynamicArray class
add_card(slide, Inches(0.5), Inches(1.4), Inches(3.8), Inches(5.2),
         title="📦 class DynamicArray",
         content_lines=[
             "Thuộc tính:",
             "  • data: Array",
             "  • length: number",
             "  • capacity: number",
             "",
             "Phương thức:",
             "  • push() / pop()",
             "  • get() / set()",
             "  • insertAt() / removeAt()",
             "  • find() / findIndex() / filter()",
             "  • sort() → QuickSort",
             "  • binarySearch()",
             "  • map() / forEach() / reduce()",
             "  • clone() / toArray()",
         ],
         border_color=PRIMARY, title_color=PRIMARY)

# Student class
add_card(slide, Inches(4.6), Inches(1.4), Inches(3.9), Inches(5.2),
         title="👤 class Student",
         content_lines=[
             "Thuộc tính:",
             "  • id, name, dob, gender",
             "  • className, email, phone",
             "  • address, scores{}",
             "  • createdAt, updatedAt",
             "",
             "Phương thức:",
             "  • calculateGPA()",
             "  • getRank()",
             "  • getAge() / getFormattedDob()",
             "  • validate()",
             "  • update() / clone()",
             "  • toJSON() / fromJSON()",
         ],
         border_color=ACCENT, title_color=ACCENT)

# StudentManager class
add_card(slide, Inches(8.8), Inches(1.4), Inches(4.0), Inches(5.2),
         title="📊 class StudentManager",
         content_lines=[
             "Thuộc tính:",
             "  • students: DynamicArray",
             "  • storageKey: string",
             "",
             "CRUD:",
             "  • addStudent() / deleteStudent()",
             "  • updateStudent()",
             "  • getStudentById()",
             "",
             "Tìm kiếm & Sắp xếp:",
             "  • searchStudents(criteria)",
             "  • sortStudents(field, order)",
             "",
             "Thống kê & Lưu trữ:",
             "  • getStatistics()",
             "  • saveToStorage()",
             "  • exportToJSON()",
         ],
         border_color=ORANGE, title_color=ORANGE)

# Arrows between classes
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.35), Inches(3.8), Inches(0.3), Inches(0.25))
arrow.fill.solid()
arrow.fill.fore_color.rgb = GRAY_TEXT
arrow.line.fill.background()

arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(3.8), Inches(0.3), Inches(0.25))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = GRAY_TEXT
arrow2.line.fill.background()


# ================================================================
# SLIDE 10: Thuật toán & Độ phức tạp tổng hợp
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "4", "BẢNG TỔNG HỢP ĐỘ PHỨC TẠP", "Phân tích thuật toán trong hệ thống")
add_bottom_bar(slide, "Phần 4: Thuật toán & Độ phức tạp", PRIMARY)
add_slide_number(slide, 10, TOTAL_SLIDES)

# Main complexity table
complex_data = [
    ["Thao tác", "Thuật toán", "Thời gian", "Không gian", "Ghi chú"],
    ["push()", "Amortized", "O(1)*", "O(1)", "Nhân đôi capacity khi đầy"],
    ["pop() / get() / set()", "Direct access", "O(1)", "O(1)", "Truy cập ngẫu nhiên"],
    ["insertAt() / removeAt()", "Shift elements", "O(n)", "O(1)", "Dịch chuyển phần tử"],
    ["sort()", "QuickSort", "O(n log n)*", "O(log n)", "Divide & Conquer"],
    ["binarySearch()", "Binary Search", "O(log n)", "O(1)", "Yêu cầu sorted"],
    ["searchStudents()", "Linear Search", "O(n)", "O(k)", "k = số kết quả"],
    ["getStatistics()", "Single-pass", "O(n)", "O(c)", "c = số lớp, 1 lần duyệt"],
    ["addStudent()", "Check + Push", "O(n)", "O(1)", "Kiểm tra trùng mã SV"],
    ["sortStudents()", "Clone + QuickSort", "O(n log n)", "O(n)", "Clone mảng + sort"],
]

table = slide.shapes.add_table(len(complex_data), 5, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.2)).table
table.columns[0].width = Inches(2.8)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(1.8)
table.columns[3].width = Inches(1.8)
table.columns[4].width = Inches(3.9)

for row_idx, row_data in enumerate(complex_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        if row_idx == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY
        else:
            p.font.color.rgb = DARK_TEXT
            if col_idx == 2:
                p.font.bold = True
                p.font.color.rgb = ACCENT
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

# Key insights
insights = [
    ("QuickSort", "O(n log n) trung bình — nhanh gấp nhiều lần BubbleSort O(n²)", ACCENT),
    ("Binary Search", "O(log n) — tìm trong 1000 SV chỉ cần ~10 bước so với 1000 bước Linear", BLUE),
    ("Single-pass Stats", "O(n) — tính tất cả thống kê chỉ trong 1 lần duyệt mảng", ORANGE),
]

for i, (title, desc, color) in enumerate(insights):
    left = Inches(0.5 + i * 4.2)
    add_card(slide, left, Inches(5.0), Inches(3.9), Inches(1.7),
             title=f"⚡ {title}", content=desc, title_color=color, border_color=color)


# ================================================================
# SLIDE 11: Giao diện ứng dụng
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "5", "GIAO DIỆN ỨNG DỤNG", "Thiết kế UI/UX hiện đại, thân thiện người dùng")
add_bottom_bar(slide, "Phần 5: Demo ứng dụng", PRIMARY)
add_slide_number(slide, 11, TOTAL_SLIDES)

# Layout description
add_card(slide, Inches(0.6), Inches(1.4), Inches(6.0), Inches(2.5),
         title="🎨 Thiết kế giao diện",
         content_lines=[
             "• Thiết kế giao diện sáng (Light Theme)",
             "• Phối màu Indigo (#4F46E5) chuyên nghiệp",
             "• Font Inter (Google Fonts) - dễ đọc",
             "• Font Awesome 6.4.0 cho icons",
             "• Responsive - tương thích mọi màn hình",
             "• Animation mượt mà, trực quan",
         ],
         border_color=PRIMARY)

# Layout mockup
layout_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.4), Inches(5.7), Inches(2.5)
)
layout_box.fill.solid()
layout_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
layout_box.line.color.rgb = RGBColor(203, 213, 225)

# Mockup header
header_mock = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.4)
)
header_mock.fill.solid()
header_mock.fill.fore_color.rgb = PRIMARY
header_mock.line.fill.background()
tf = header_mock.text_frame
p = tf.paragraphs[0]
p.text = "🎓 HỆ THỐNG QUẢN LÝ SINH VIÊN"
p.font.color.rgb = WHITE
p.font.size = Pt(10)
p.alignment = PP_ALIGN.CENTER

# Mockup sidebar
sidebar_mock = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(1.95), Inches(1.3), Inches(1.8)
)
sidebar_mock.fill.solid()
sidebar_mock.fill.fore_color.rgb = WHITE
sidebar_mock.line.color.rgb = RGBColor(226, 232, 240)
tf = sidebar_mock.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "  Tổng quan"
p.font.size = Pt(8)
p.font.color.rgb = PRIMARY
for item in ["  Thêm SV", "  Danh sách", "  Tìm kiếm", "  Thống kê"]:
    pp = tf.add_paragraph()
    pp.text = item
    pp.font.size = Pt(8)
    pp.font.color.rgb = GRAY_TEXT

# Mockup content area
content_mock = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(1.95), Inches(4.1), Inches(1.8)
)
content_mock.fill.solid()
content_mock.fill.fore_color.rgb = WHITE
content_mock.line.color.rgb = RGBColor(226, 232, 240)
tf = content_mock.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "     [Dashboard / Nội dung chính]"
p.font.size = Pt(9)
p.font.color.rgb = GRAY_TEXT
p.alignment = PP_ALIGN.CENTER

# Feature screens
screens = [
    ("📊 Tổng quan", "Dashboard với thẻ thống kê\nvà SV mới thêm gần đây"),
    ("➕ Thêm sinh viên", "Form nhập liệu với\nvalidation đầy đủ"),
    ("📋 Danh sách SV", "Bảng hiển thị phân trang\nSắp xếp, Xem, Sửa, Xóa"),
    ("🔍 Tìm kiếm", "Tìm theo từ khóa\n+ bộ lọc đa tiêu chí"),
    ("📈 Thống kê", "Biểu đồ tròn & cột\nChart.js trực quan"),
]

for i, (title, desc) in enumerate(screens):
    left = Inches(0.6 + i * 2.5)
    add_card(slide, left, Inches(4.3), Inches(2.3), Inches(2.2),
             title=title, content=desc,
             title_color=PRIMARY if i % 2 == 0 else ACCENT)


# ================================================================
# SLIDE 12: Demo chức năng
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "5", "CHỨC NĂNG CHI TIẾT", "Các chức năng chính của hệ thống")
add_bottom_bar(slide, "Phần 5: Demo ứng dụng", PRIMARY)
add_slide_number(slide, 12, TOTAL_SLIDES)

features = [
    ("CRUD Operations", "📝", PRIMARY,
     ["✓ Thêm sinh viên (validation đầy đủ)",
      "✓ Xem chi tiết qua modal popup",
      "✓ Sửa thông tin inline",
      "✓ Xóa với xác nhận"]),
    ("Tìm kiếm", "🔍", ACCENT,
     ["✓ Tìm theo mã SV, họ tên, lớp",
      "✓ Lọc theo xếp loại học lực",
      "✓ Lọc theo giới tính",
      "✓ Kết hợp nhiều tiêu chí"]),
    ("Sắp xếp", "📊", BLUE,
     ["✓ Theo mã SV (A-Z, Z-A)",
      "✓ Theo họ tên (A-Z, Z-A)",
      "✓ Theo GPA (cao→thấp, thấp→cao)",
      "✓ Theo lớp"]),
    ("Thống kê & Biểu đồ", "📈", ORANGE,
     ["✓ Tổng số SV, điểm TB, cao/thấp nhất",
      "✓ Biểu đồ tròn: phân bố xếp loại",
      "✓ Biểu đồ cột: thống kê theo lớp",
      "✓ Bảng chi tiết từng lớp"]),
]

for i, (title, icon, color, items) in enumerate(features):
    row = i // 2
    col = i % 2
    left = Inches(0.6 + col * 6.3)
    top = Inches(1.4 + row * 2.8)

    add_card(slide, left, top, Inches(5.9), Inches(2.5),
             title=f"{icon} {title}",
             content_lines=items,
             title_color=color, border_color=color)

# Extra features at bottom
extras = [
    "💾 Lưu trữ LocalStorage tự động",
    "📤 Xuất dữ liệu JSON",
    "📥 Import dữ liệu từ JSON",
    "🎲 Dữ liệu mẫu 30 sinh viên",
    "⏰ Đồng hồ realtime",
]

for i, text in enumerate(extras):
    left = Inches(0.6 + i * 2.5)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(6.2), Inches(2.3), Inches(0.55)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = RGBColor(226, 232, 240)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER


# ================================================================
# SLIDE 13: Kết quả & Kết luận
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "6", "KẾT QUẢ & KẾT LUẬN", "Tổng kết các kết quả đạt được")
add_bottom_bar(slide, "Phần 6: Kết luận", PRIMARY)
add_slide_number(slide, 13, TOTAL_SLIDES)

# Results
results = [
    ("✅ Mảng động", "Xây dựng thành công cấu trúc\nDynamic Array với đầy đủ\n20+ phương thức"),
    ("✅ QuickSort", "Cài đặt thuật toán QuickSort\nĐộ phức tạp O(n log n)\nSắp xếp đa tiêu chí"),
    ("✅ Tìm kiếm", "Linear Search + Binary Search\nTìm kiếm đa tiêu chí\nLọc dữ liệu linh hoạt"),
    ("✅ Giao diện", "UI hiện đại, responsive\nChart.js biểu đồ\n5 màn hình chức năng"),
    ("✅ OOP", "Code tổ chức theo OOP\n9 file JavaScript\nCấu trúc rõ ràng"),
    ("✅ Lưu trữ", "LocalStorage persistence\nExport/Import JSON\nDữ liệu mẫu 30 SV"),
]

for i, (title, desc) in enumerate(results):
    row = i // 3
    col = i % 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.45 + row * 2.3)
    add_card(slide, left, top, Inches(3.8), Inches(2.0),
             title=title, content=desc,
             title_color=ACCENT, border_color=ACCENT)

# Limitations & Future
add_card(slide, Inches(0.6), Inches(5.3), Inches(5.8), Inches(1.4),
         title="⚠️ Hạn chế",
         content_lines=[
             "• Dữ liệu chỉ lưu cục bộ (LocalStorage ~5MB)",
             "• Chưa có chức năng đăng nhập, phân quyền",
             "• Chưa tích hợp database server",
         ],
         border_color=YELLOW, title_color=YELLOW)

add_card(slide, Inches(6.9), Inches(5.3), Inches(5.8), Inches(1.4),
         title="🚀 Hướng phát triển",
         content_lines=[
             "• Tích hợp backend (Node.js / Python)",
             "• Sử dụng database (MongoDB, MySQL)",
             "• Thêm import/export Excel, phân quyền",
         ],
         border_color=BLUE, title_color=BLUE)


# ================================================================
# SLIDE 14: Thank you
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, PRIMARY, PRIMARY_DARK)

# Thank you text
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CẢM ƠN THẦY/CÔ ĐÃ LẮNG NGHE!"
p.font.color.rgb = WHITE
p.font.size = Pt(40)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Accent line
accent = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.2), Inches(3.333), Inches(0.06)
)
accent.fill.solid()
accent.fill.fore_color.rgb = ACCENT
accent.line.fill.background()

# Project summary
txBox = slide.shapes.add_textbox(Inches(2), Inches(3.6), Inches(9.333), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Đề tài: Xây dựng phần mềm Quản lý Sinh viên"
p.font.color.rgb = RGBColor(199, 210, 254)
p.font.size = Pt(18)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Sử dụng cấu trúc Mảng Động & Giao diện Đồ họa"
p2.font.color.rgb = RGBColor(165, 180, 252)
p2.font.size = Pt(15)
p2.alignment = PP_ALIGN.CENTER

# Q&A
txBox = slide.shapes.add_textbox(Inches(2), Inches(5.0), Inches(9.333), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "❓ Sẵn sàng trả lời câu hỏi"
p.font.color.rgb = WHITE
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER

# Student info box
info_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(5.8), Inches(5.333), Inches(0.8)
)
info_box.fill.solid()
info_box.fill.fore_color.rgb = RGBColor(67, 56, 202)
info_box.line.color.rgb = SECONDARY
info_box.line.width = Pt(1)
tf = info_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Sinh viên: [Điền tên]  •  GVHD: [Điền tên]"
p.font.color.rgb = RGBColor(199, 210, 254)
p.font.size = Pt(13)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Môn: Lập trình nâng cao  •  Năm học: 2025-2026"
p2.font.color.rgb = RGBColor(165, 180, 252)
p2.font.size = Pt(11)
p2.alignment = PP_ALIGN.CENTER


# ================================================================
# SAVE FILE
# ================================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "BaoVe_BTL.pptx")
prs.save(output_path)
print(f"✅ Đã tạo thành công file: {output_path}")
print(f"📄 Tổng số slide: {len(prs.slides)}")
print(f"📐 Kích thước: Widescreen 16:9")
